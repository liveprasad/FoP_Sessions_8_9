"""
Generate Session 8 and Session 9 PowerPoint files for FoP (Fundamentals of Programming).
Educational theme: soft background, navy titles, readable body text.
Upload the .pptx to Google Slides for editing/presenting.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUTPUT_DIR = "."

# Educational theme (calm, readable, classroom-friendly)
THEME_BG = RGBColor(245, 248, 252)       # cool off-white / light blue-gray
THEME_TITLE = RGBColor(21, 67, 108)    # navy
THEME_SUBTITLE = RGBColor(70, 90, 110) # muted slate
THEME_BODY = RGBColor(44, 44, 46)      # near-black for bullets
def _fill_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_title_shape(shape, font_size_pt=30, subtitle=False):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.space_after = Pt(6)
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(font_size_pt - 4) if subtitle else Pt(font_size_pt)
            run.font.bold = not subtitle
            run.font.color.rgb = THEME_SUBTITLE if subtitle else THEME_TITLE


def _style_body_text_frame(tf, font_size_pt=18):
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for para in tf.paragraphs:
        para.space_after = Pt(10)
        para.line_spacing = 1.15
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = THEME_BODY


def _apply_theme_to_slide(slide, prs_height, is_title_slide=False):
    _fill_slide_background(slide, THEME_BG)
    if slide.shapes.title:
        _style_title_shape(slide.shapes.title, 32 if is_title_slide else 28, subtitle=False)
    if is_title_slide and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        if ph.has_text_frame:
            _style_title_shape(ph, 20, subtitle=True)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    _apply_theme_to_slide(slide, prs.slide_height, is_title_slide=True)


def add_content_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _fill_slide_background(slide, THEME_BG)
    slide.shapes.title.text = title
    _style_title_shape(slide.shapes.title, 26, subtitle=False)

    body = slide.placeholders[1].text_frame
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = THEME_BODY

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide


def add_ref_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _fill_slide_background(slide, THEME_BG)
    slide.shapes.title.text = "References"
    _style_title_shape(slide.shapes.title, 26, subtitle=False)

    body = slide.placeholders[1].text_frame
    refs = [
        "Downey, A. (2012). Think Python. O'Reilly Media, Inc.",
        "Matthes, E. (2023). Python Crash Course: A hands-on, project-based introduction to programming.",
        "Müller, A. C., & Guido, S. (2016). Introduction to machine learning with Python: a guide for data scientists. O'Reilly Media, Inc.",
        "GeeksforGeeks Python: https://www.geeksforgeeks.org/python/",
        "W3Schools Python: https://www.w3schools.com/python/",
    ]
    for i, r in enumerate(refs):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = r
        p.level = 0
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(15)
            run.font.color.rgb = THEME_BODY


def build_session8_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs,
        "Session 8: Tuples, Sets & Dictionaries",
        "FoP – B.S. in Management and Public Policy | 90 min"
    )
    add_content_slide(prs, "Learning objectives", [
        "Organize data using tuples, sets, and dictionaries.",
        "Choose the right structure for policy and management data (schemes, beneficiaries, surveys).",
        "Apply to real-world tasks: unique stakeholders, scheme lookups, survey response counts.",
    ])
    add_content_slide(prs, "Tuples – management & policy use", [
        "Immutable ordered sequence: (scheme_name, year, budget) — e.g. one fixed record.",
        "Use when: order matters and data should not change (e.g. policy snapshot, KPI row).",
        "Unpacking: scheme_name, year, budget = policy_record.",
    ], notes=(
        "INTERNALS (for presenter):\n"
        "• Tuples are immutable: once created, you cannot add/remove/change elements. Python can optimize storage and reuse.\n"
        "• Stored as a fixed-size sequence in memory; indexing by position is O(1).\n"
        "• Tuples are hashable if all elements are hashable — so they can be used as dict keys or set elements (e.g. (scheme_id, year)).\n"
        "• Less memory overhead than lists when the sequence never changes."
    ))
    add_content_slide(prs, "Sets – unique entities", [
        "Unordered collection of unique elements: e.g. unique beneficiary IDs, districts, departments.",
        "Use when: counting distinct stakeholders, regions covered, or removing duplicate responses.",
        "Operations: union (combined coverage), intersection (common beneficiaries), difference.",
    ], notes=(
        "INTERNALS (for presenter):\n"
        "• Sets are implemented with a hash table: each element must be hashable (no lists/dicts inside the set).\n"
        "• Membership check (x in s) is O(1) on average — very fast for large collections.\n"
        "• Order is not guaranteed; Python may reorder for efficiency. Do not rely on order.\n"
        "• Adding/removing is O(1) average. Union, intersection, difference are efficient."
    ))
    add_content_slide(prs, "Dictionaries – key–value lookups", [
        "Key–value pairs: scheme_id → details, department → head, indicator → target.",
        "Use when: looking up policy details, program metadata, or survey codes by ID.",
        ".keys(), .values(), .items() for reporting and iteration.",
    ], notes=(
        "INTERNALS (for presenter):\n"
        "• Dicts are implemented with a hash table: keys must be hashable (e.g. str, int, tuple of hashables).\n"
        "• Lookup by key (d[key]) is O(1) on average — ideal for scheme_id → details.\n"
        "• Insertion and deletion by key are also O(1) average.\n"
        "• From Python 3.7+ dicts preserve insertion order; iterating is in the order keys were added."
    ))
    add_content_slide(prs, "Under the hood: Tuples, Sets, Dictionaries", [
        "Tuple: immutable, fixed sequence; hashable if elements are hashable; O(1) index access.",
        "Set: hash table; O(1) membership; elements must be hashable; order not guaranteed.",
        "Dict: hash table; O(1) lookup by key; keys must be hashable; insertion order preserved (3.7+).",
    ], notes=(
        "Use this slide + notes to explain internals when students ask \"why is set/dict fast?\" or \"why can't I put a list in a set?\".\n\n"
        "TUPLES: Immutability allows Python to optimize; no reallocation. Hashable tuples can go in sets/dict keys.\n\n"
        "SETS: Hash table gives O(1) 'in' check. Uniqueness is enforced by hash + equality. So only hashable types allowed.\n\n"
        "DICTIONARIES: Same idea — hash(key) determines bucket; O(1) average for get/set/del. .get(key, default) avoids KeyError."
    ))
    add_content_slide(prs, "When to use which? (management & policy)", [
        "Tuple: one fixed record (e.g. scheme name, financial year, allocated amount).",
        "Set: unique beneficiaries, unique districts in a program, unique response categories.",
        "Dict: scheme_id → full details; department → contact; indicator_code → definition.",
    ])
    add_content_slide(prs, "Practical in-class (Colab)", [
        "Session 8 Student notebook: tuples (policy record), sets (unique regions/beneficiaries), dicts (scheme lookup).",
        "Hands-on: e.g. count unique districts in a program, build a scheme lookup table.",
        "All examples use business management and public policy contexts.",
    ])
    add_content_slide(prs, "Session 8 – Recap", [
        "Tuples: immutable, ordered.",
        "Sets: unique, unordered, fast membership.",
        "Dictionaries: key–value, fast lookup by key.",
    ])
    add_ref_slide(prs)

    path = f"{OUTPUT_DIR}/Session8_Tuples_Sets_Dictionaries.pptx"
    prs.save(path)
    print(f"Saved: {path}")


def build_session9_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs,
        "Session 9: Functions – Basics",
        "FoP – B.S. in Management and Public Policy | 90 min"
    )
    add_content_slide(prs, "Learning objectives", [
        "Define and call functions with parameters and return values.",
        "Write reusable code for policy and management tasks (formatting, summaries, lookups).",
        "Apply functions: budget display, scheme summaries, indicator lookups, report snippets.",
    ])
    add_content_slide(prs, "Why functions? (management & policy)", [
        "Reuse: format budget in lakhs/crores, format scheme names, or build report lines repeatedly.",
        "Structure: break reporting into clear steps (e.g. get_scheme_summary, format_for_display).",
        "Easier to maintain when definitions change (e.g. currency, rounding rules).",
    ])
    add_content_slide(prs, "Defining a function", [
        "def function_name(parameter1, parameter2):",
        "    '''Optional docstring — e.g. \"Format budget in lakhs for display\"'''",
        "    # body",
        "    return result",
    ], notes=(
        "INTERNALS (for presenter):\n"
        "• def is a statement that creates a function object and binds it to the name. The body is not run until the function is called.\n"
        "• Docstrings (first string after def) are stored in .__doc__ and used by help().\n"
        "• Indentation defines the body; first dedent ends the function."
    ))
    add_content_slide(prs, "Parameters and return", [
        "Parameters: e.g. scheme name and year, or a list of grant amounts.",
        "return: e.g. formatted string, summary dict, or lookup result.",
        "Without return: function returns None (avoid for data you need to use).",
    ], notes=(
        "INTERNALS (for presenter):\n"
        "• Arguments are passed by object reference: the parameter name refers to the same object. Mutating (e.g. list.append) is visible to the caller; reassigning the name is not.\n"
        "• return exits the function immediately and sends the value back. Only one return is executed per call.\n"
        "• Default args are evaluated once at def time; avoid mutable defaults (e.g. def f(x=[]) )."
    ))
    add_content_slide(prs, "Scope", [
        "Variables inside a function are local (e.g. temp totals, formatted strings).",
        "Variables in the notebook are global; prefer passing data in and returning results.",
        "Keeps policy/data inputs explicit and outputs clear for reporting.",
    ], notes=(
        "INTERNALS (for presenter):\n"
        "• LEGB: Local, Enclosing (nested functions), Global, Built-in. Assignment inside a function creates a local name unless declared global.\n"
        "• Prefer passing inputs as arguments and returning results — makes dependencies clear and functions testable.\n"
        "• Reading a global is fine; modifying it inside a function requires global declaration (discourage for clarity)."
    ))
    add_content_slide(prs, "Under the hood: functions", [
        "def creates a function object; body runs only when called.",
        "Arguments passed by object reference; return exits and sends a value back.",
        "Scope: local names inside the function; prefer pass-in, return-out for clarity.",
    ], notes=(
        "Use this when students ask \"when is the body executed?\" or \"why did my list change?\".\n\n"
        "EXECUTION: def runs once and binds the function; each call runs the body with fresh local names.\n\n"
        "PASS BY REFERENCE: If you pass a list and append to it, the caller sees the change. If you do param = something_else, the caller's variable is unchanged.\n\n"
        "RETURN: Can return any type — string, dict, set, tuple, or a tuple of (set, dict) for multiple results."
    ))
    add_content_slide(prs, "Practical in-class (Colab)", [
        "Session 9 Student notebook: functions for policy/scheme formatting, grant averages, scheme lookup.",
        "Hands-on: e.g. format_budget(amount), summarize_grants(list), get_scheme_by_id(dict, id).",
        "All examples aligned with business management and public policy.",
    ])
    add_content_slide(prs, "Session 9 – Recap", [
        "def name(params): ... return value",
        "Use functions for reusable formatting, summaries, and lookups in policy/management data.",
        "Structure programs with small, named steps for clarity and maintenance.",
    ])
    add_ref_slide(prs)

    path = f"{OUTPUT_DIR}/Session9_Functions_Basics.pptx"
    prs.save(path)
    print(f"Saved: {path}")


if __name__ == "__main__":
    build_session8_pptx()
    build_session9_pptx()
    print("Done. Upload .pptx files to Google Slides as needed.")
